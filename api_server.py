#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT to Images API Server

提供 HTTP API 服务，接收 PPT 文件并返回图片 URL
运行在 localhost:4000

支持异步转换和进度查询
"""

from fastapi import FastAPI, UploadFile, File, HTTPException, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import HTMLResponse, JSONResponse
import os
import uuid
import tempfile
import time
import threading
import glob
from pathlib import Path
from typing import List, Dict, Any, Optional, Callable
from dataclasses import dataclass, asdict
from enum import Enum
from datetime import datetime
from ppt_exporter import PPTExporter, ExportConfig

# 尝试导入 python-pptx 用于获取幻灯片数量
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

app = FastAPI(title="PPT to Images API", version="2.0.0")

# CORS 设置
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # 允许所有来源（开发环境）
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 输出目录
OUTPUT_BASE_DIR = os.path.join(os.path.dirname(__file__), "output")
os.makedirs(OUTPUT_BASE_DIR, exist_ok=True)

# 挂载静态文件目录
app.mount("/images", StaticFiles(directory=OUTPUT_BASE_DIR), name="images")

# HTML 模板目录
TEMPLATES_DIR = os.path.join(os.path.dirname(__file__), "templates")
os.makedirs(TEMPLATES_DIR, exist_ok=True)

# API 基础 URL（支持环境变量配置）
API_BASE_URL = os.getenv("API_BASE_URL", "http://localhost:4000")

# 任务状态枚举
class TaskStatus(str, Enum):
    PENDING = "pending"      # 等待处理
    PROCESSING = "processing"  # 处理中
    COMPLETED = "completed"   # 完成
    FAILED = "failed"        # 失败

# 任务信息数据类
@dataclass
class TaskInfo:
    task_id: str  # task_id 同时也是 folder_id
    filename: str
    status: TaskStatus
    progress: int  # 0-100
    total_slides: int
    current_slide: int
    created_at: str
    updated_at: str
    status_message: str = ""  # 当前状态描述
    images: List[Dict[str, Any]] = None
    error: Optional[str] = None
    
    def to_dict(self):
        return asdict(self)

# 内存中的任务缓存
tasks_cache: Dict[str, TaskInfo] = {}


def get_ppt_slide_count(ppt_path: str) -> int:
    """获取 PPT 的幻灯片数量（只计算非隐藏幻灯片）"""
    if not HAS_PPTX:
        return 0
        
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(ppt_path)
        
        # 尝试检测隐藏的幻灯片
        # 注意：python-pptx 没有直接的 API 来检查幻灯片是否隐藏
        # 所以我们假设所有幻灯片都是可见的
        visible_count = len(prs.slides)
        
        # 检查是否有特殊标记的幻灯片
        hidden_count = 0
        for i, slide in enumerate(prs.slides, 1):
            # 检查 slide 的 XML 中是否有 show="0" 属性
            try:
                slide_xml = slide.element.xml
                if b'show="0"' in slide_xml or b'show="false"' in slide_xml:
                    hidden_count += 1
                    print(f"  📌 幻灯片 {i} 是隐藏的")
            except:
                pass
        
        visible_count = len(prs.slides) - hidden_count
        
        if hidden_count > 0:
            print(f"  📊 总共 {len(prs.slides)} 张幻灯片，其中 {hidden_count} 张隐藏，可见 {visible_count} 张")
        else:
            print(f"  📊 总共 {visible_count} 张幻灯片（无隐藏幻灯片）")
            
        return visible_count
        
    except Exception as e:
        print(f"⚠️ 无法读取 PPT 页数: {e}")
        return 0


def progress_callback(task_id: str, current: int, total: int, filename: str):
    """进度回调函数"""
    task = tasks_cache.get(task_id)
    if not task:
        return
    
    task.current_slide = current
    task.total_slides = total
    
    # 计算进度 (30% 到 85% 之间)
    progress = 30 + int((current / total) * 55)
    task.progress = min(progress, 85)
    
    task.status_message = f"✓ 已保存: {filename} ({current}/{total})"
    task.updated_at = datetime.now().isoformat()
    
    print(f"  ✓ 已保存: {filename} ({current}/{total})")


@app.get("/", response_class=HTMLResponse)
async def root():
    """返回 Web 界面"""
    html_file = os.path.join(TEMPLATES_DIR, "index.html")
    
    if os.path.exists(html_file):
        with open(html_file, 'r', encoding='utf-8') as f:
            return HTMLResponse(content=f.read())
    else:
        return HTMLResponse(
            content="<h1>Error: Template not found</h1><p>Please ensure templates/index.html exists.</p>",
            status_code=500
        )


@app.get("/health")
def health_check():
    return {
        "status": "healthy",
        "service": "ppt-to-images",
        "port": 4000,
        "tasks_count": len(tasks_cache)
    }


@app.post("/api/convert-async")
async def convert_ppt_async(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    dpi: int = 300,
    format: str = "png"
):
    """
    异步转换 PPT 为图片
    立即返回任务 ID，后台处理
    """
    
    # 验证文件类型
    if not file.filename.endswith(('.ppt', '.pptx')):
        raise HTTPException(
            status_code=400,
            detail="仅支持 .ppt 或 .pptx 格式"
        )
    
    # 生成任务 ID（同时作为 folder_id）
    task_id = str(uuid.uuid4())
    
    # 创建任务记录
    task_info = TaskInfo(
        task_id=task_id,  # task_id 就是 folder_id
        filename=file.filename,
        status=TaskStatus.PENDING,
        progress=0,
        total_slides=0,
        current_slide=0,
        status_message="等待处理...",
        created_at=datetime.now().isoformat(),
        updated_at=datetime.now().isoformat(),
        images=[]
    )
    
    tasks_cache[task_id] = task_info
    
    # 保存文件
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    temp_file_path = temp_file.name
    
    contents = await file.read()
    temp_file.write(contents)
    temp_file.close()
    
    # 添加后台任务
    background_tasks.add_task(
        process_ppt_task,
        task_id=task_id,
        temp_file_path=temp_file_path,
        dpi=dpi,
        format=format
    )
    
    return {
        "success": True,
        "task_id": task_id,
        "message": "任务已创建，正在处理中"
    }


def process_ppt_task(task_id: str, temp_file_path: str, dpi: int, format: str):
    """后台处理 PPT 转换任务"""
    
    task = tasks_cache.get(task_id)
    if not task:
        return
    
    monitor_thread = None
    
    try:
        # 阶段 1: 准备转换 - 先读取幻灯片数量
        task.status = TaskStatus.PROCESSING
        task.progress = 5
        task.status_message = "读取 PPT 信息..."
        task.updated_at = datetime.now().isoformat()
        
        # 获取总页数
        slide_count = get_ppt_slide_count(temp_file_path)
        if slide_count > 0:
            task.total_slides = slide_count
            task.status_message = f"检测到 {slide_count} 张幻灯片，准备转换..."
            print(f"📊 检测到 {slide_count} 张幻灯片")
        else:
            task.status_message = "准备转换..."
            print("📊 开始转换...")
        task.updated_at = datetime.now().isoformat()
        time.sleep(0.5)
        
        # 使用 task_id 作为输出文件夹名
        output_dir = os.path.join(OUTPUT_BASE_DIR, task_id)
        os.makedirs(output_dir, exist_ok=True)
        
        # 阶段 2: 开始转换为 PDF
        task.progress = 10
        task.status_message = "转换为 PDF..."
        task.updated_at = datetime.now().isoformat()
        print("📄 转换为 PDF...")
        time.sleep(0.3)
        
        # 创建导出器
        config = ExportConfig(dpi=dpi, format=format, quality=95)
        exporter = PPTExporter(config)
        
        # 阶段 3: 使用 LibreOffice 转换 PDF
        task.progress = 15
        task.status_message = "正在使用 LibreOffice 转换 PDF..."
        task.updated_at = datetime.now().isoformat()
        print("⚙️  正在使用 LibreOffice 转换...")
        
        # 阶段 4: PDF 转换为图片（这里是实际的转换过程）
        task.progress = 30
        if task.total_slides > 0:
            task.status_message = f"PDF 转换为图片 (0/{task.total_slides})..."
        else:
            task.status_message = "PDF 转换为图片..."
        task.updated_at = datetime.now().isoformat()
        print("🖼️  开始转换为图片...")
        print(f"   输出目录: {output_dir}")
        print(f"   文件格式: {format}")
        
        # 确保输出目录存在
        os.makedirs(output_dir, exist_ok=True)
        
        # 导出图片（使用回调函数报告进度）
        print("⚙️  调用 PPTExporter.export()...")
        
        # 创建进度回调
        def on_progress(current, total, filename):
            progress_callback(task_id, current, total, filename)
        
        image_files = exporter.export(
            ppt_path=temp_file_path,
            output_dir=output_dir,
            method="auto",
            prefix="slide",
            progress_callback=on_progress
        )
        print(f"✅ PPTExporter.export() 完成，返回 {len(image_files)} 个文件")
        
        # 更新为实际的图片数量
        actual_count = len(image_files)
        if actual_count != task.total_slides:
            print(f"⚠️  预期 {task.total_slides} 张，实际生成 {actual_count} 张")
            task.total_slides = actual_count
        
        task.current_slide = actual_count
        
        # 阶段 5: 生成图片 URL
        task.progress = 90
        task.status_message = f"正在生成图片 URL ({actual_count} 张)..."
        task.updated_at = datetime.now().isoformat()
        print(f"🔗 生成 URL ({actual_count} 张)...")
        time.sleep(0.3)
        
        # 生成图片 URL（使用配置的 API_BASE_URL）
        image_urls = [
            {
                "slide_number": i + 1,
                "url": f"{API_BASE_URL}/images/{task_id}/{os.path.basename(img_path)}",
                "filename": os.path.basename(img_path)
            }
            for i, img_path in enumerate(image_files)
        ]
        
        # 更新任务状态
        task.status = TaskStatus.COMPLETED
        task.progress = 100
        task.current_slide = actual_count
        task.status_message = f"转换完成！共 {actual_count} 张图片"
        task.images = image_urls
        task.updated_at = datetime.now().isoformat()
        
    except Exception as e:
        # 失败处理
        task.status = TaskStatus.FAILED
        task.status_message = "转换失败"
        task.error = str(e)
        task.updated_at = datetime.now().isoformat()
        print(f"❌ 转换失败: {e}")
        
        # 清理失败的文件夹（task_id 就是文件夹名）
        folder_path = os.path.join(OUTPUT_BASE_DIR, task_id)
        if os.path.exists(folder_path):
            import shutil
            shutil.rmtree(folder_path, ignore_errors=True)
    
    finally:
        # 清理临时文件
        try:
            os.unlink(temp_file_path)
        except:
            pass


@app.get("/api/task/{task_id}")
def get_task_status(task_id: str):
    """查询任务状态"""
    
    task = tasks_cache.get(task_id)
    
    if not task:
        raise HTTPException(
            status_code=404,
            detail="任务不存在"
        )
    
    return {
        "success": True,
        "task": task.to_dict()
    }


@app.get("/api/tasks")
def list_all_tasks():
    """列出所有任务"""
    
    tasks = [task.to_dict() for task in tasks_cache.values()]
    
    # 按创建时间倒序排序
    tasks.sort(key=lambda x: x['created_at'], reverse=True)
    
    return {
        "success": True,
        "count": len(tasks),
        "tasks": tasks
    }


@app.delete("/api/task/{task_id}")
def delete_task(task_id: str):
    """删除任务"""
    
    task = tasks_cache.get(task_id)
    
    if not task:
        raise HTTPException(
            status_code=404,
            detail="任务不存在"
        )
    
    # 删除文件夹（task_id 就是文件夹名）
    folder_path = os.path.join(OUTPUT_BASE_DIR, task_id)
    if os.path.exists(folder_path):
        import shutil
        shutil.rmtree(folder_path, ignore_errors=True)
    
    # 删除任务记录
    del tasks_cache[task_id]
    
    return {
        "success": True,
        "message": f"任务 {task_id} 已删除"
    }


# 保持原有的同步 API（向后兼容）
@app.post("/api/convert")
async def convert_ppt_sync(
    file: UploadFile = File(...),
    dpi: int = 300,
    format: str = "png"
):
    """
    同步转换 PPT 为图片（向后兼容）
    """
    
    if not file.filename.endswith(('.ppt', '.pptx')):
        raise HTTPException(
            status_code=400,
            detail="仅支持 .ppt 或 .pptx 格式文件"
        )
    
    try:
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        temp_file_path = temp_file.name
        
        contents = await file.read()
        temp_file.write(contents)
        temp_file.close()
        
        # 生成 folder_id（使用 UUID）
        folder_id = str(uuid.uuid4())
        output_dir = os.path.join(OUTPUT_BASE_DIR, folder_id)
        os.makedirs(output_dir, exist_ok=True)
        
        config = ExportConfig(dpi=dpi, format=format, quality=95)
        exporter = PPTExporter(config)
        
        image_files = exporter.export(
            ppt_path=temp_file_path,
            output_dir=output_dir,
            method="auto",
            prefix="slide"
        )
        
        image_urls = [
            {
                "slide_number": i + 1,
                "url": f"{API_BASE_URL}/images/{folder_id}/{os.path.basename(img_path)}",
                "filename": os.path.basename(img_path)
            }
            for i, img_path in enumerate(image_files)
        ]
        
        os.unlink(temp_file_path)
        
        return {
            "success": True,
            "folder_id": folder_id,  # folder_id 在同步模式下是独立的 UUID
            "count": len(image_urls),
            "images": image_urls,
            "message": f"成功转换 {len(image_urls)} 张图片"
        }
    
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"转换失败: {str(e)}"
        )


if __name__ == "__main__":
    import uvicorn
    
    print("=" * 60)
    print("🚀 PPT to Images API Server v2.0")
    print("=" * 60)
    print(f"📍 Web 界面: http://localhost:4000")
    print(f"📍 API 地址: http://localhost:4000")
    print(f"📚 API 文档: http://localhost:4000/docs")
    print(f"🔍 健康检查: http://localhost:4000/health")
    print("=" * 60)
    print()
    print("功能:")
    print("  • Web 界面上传和查看")
    print("  • 异步任务处理")
    print("  • 实时进度查询")
    print("  • 自动状态更新")
    print("=" * 60)
    print()
    
    uvicorn.run(
        app,
        host="0.0.0.0",
        port=4000,
        log_level="info"
    )
